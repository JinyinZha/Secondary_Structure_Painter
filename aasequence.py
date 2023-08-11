# -*- coding: utf-8 -*-
"""
Created on Mon Mar 28 14:49:58 2022

@author: 15308
"""

from utils import cm2emu

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE,MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
from Protein import Protein
my_dir = os.getcwd()

class AASequence():
    def __init__(self,paras):
        self.paras = paras
        self.file_name = paras["pptx"]
        self.fig_width = float(paras["fig_width"])
        self.char_size = int(paras["char_size"])
        self.char_length = self.char_size * 0.0216667
        self.char_height = self.char_size * 0.0433333
        self.shape_height = paras["kh"] * self.char_height
        self.text_dy = 0.5* (paras["kh"]-1) * self.char_height
        self.helix_color = RGBColor(int(paras["helix_color_r"]),int(paras["helix_color_g"]),int(paras["helix_color_b"]))
        self.sheet_color = RGBColor(int(paras["sheet_color_r"]),int(paras["sheet_color_g"]),int(paras["sheet_color_b"]))
        self.loop_color = RGBColor(int(paras["loop_color_r"]),int(paras["loop_color_g"]),int(paras["loop_color_b"]))
        self.protein = Protein(paras["protein"])
        self.error = self.protein.error
        
    def make_PPT(self):
        self.prs = Presentation(pptx="%s/default.pptx"%(my_dir))
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

    def make_rows(self):
        #1, Initialize
        res_dict = self.protein.residue_dict
        resids = list(res_dict.keys())
        resids.sort()
        #2, make rows
        rows = Rows(res_dict[resids[0]][1])
        for resi in resids:
            resname = res_dict[resi][0]
            ss = res_dict[resi][1]
            #2.1 Check row Length. If Exceed, Update.
            if rows.current_row.get_length(self.char_length) > self.fig_width:
                if ss == rows.current_row.current_shape.ss:
                    rows.current_row.current_shape.ed_id = False
                    rows.update(ss)
                    rows.current_row.current_shape.st_id = False
                else:
                    rows.update(ss)
            #2.2 Check SS. If Changed, Update.
            if ss != rows.current_row.current_shape.ss:
                rows.current_row.update(ss)
            #2.3 Write Information
            rows.current_row.current_shape.update(resname=resname,resid=str(resi),paras=self.paras,char_length=self.char_length)
        self.rows = rows
               
    def draw(self):
        sty = 0
        for row in self.rows:
            stx = 0
            line_st = []
            line_ed = []
            for sh in row:
                rotation = 0
                #Draw Shapes and Residues
                if sh.ss == "":
                    the_shape = MSO_SHAPE.OVAL
                    self.draw_shape(the_shape,cm2emu(stx),cm2emu(sty),
                                    cm2emu(sh.length),cm2emu(self.shape_height),
                                    fill_color=self.loop_color)
                    line_ed.append([stx,sty+0.5*self.shape_height])
                    line_st.append([stx+sh.length,sty+0.5*self.shape_height])
                elif sh.ss == "S":
                    the_shape = MSO_SHAPE.PENTAGON
                    self.draw_shape(the_shape,cm2emu(stx),cm2emu(sty),
                                    cm2emu(sh.length),cm2emu(self.shape_height),
                                    [(0,0.2)],fill_color=self.sheet_color)
                    line_ed.append([stx,sty+0.5*self.shape_height])
                    line_st.append([stx+sh.length,sty+0.5*self.shape_height])
                    self.draw_resids(sh, stx, sty)
                elif sh.ss == "H":
                    the_shape = MSO_SHAPE.CAN
                    rotation = 270
                    l = sh.length
                    w = self.shape_height
                    stxx = stx + 0.5 * (l - w)
                    styy = sty - 0.5 * (l - w)
                    self.draw_shape(the_shape,cm2emu(stxx),cm2emu(styy),
                                    cm2emu(w),cm2emu(l),
                                    rotation=rotation,fill_color=self.helix_color)
                    line_ed.append([stx + 0.35*self.paras["margin_can"],sty+0.5*self.shape_height])
                    line_st.append([stx+sh.length,sty+0.5*self.shape_height])
                    self.draw_resids(sh, stx, sty)
                    
                    
                    
                self.draw_textbox(cm2emu(stx),cm2emu(sty+self.text_dy),
                                  cm2emu(sh.length),cm2emu(self.shape_height),sh.margin,
                                         sh.resnames,sh.colors)
                
                stx += sh.length + self.char_length
            sty += self.shape_height * (2 + self.paras["row_gap"])
        
            for i in range(len(line_st)-1):
                self.draw_line(cm2emu(line_st[i][0]),cm2emu(line_st[i][1]),
                               cm2emu(line_ed[i+1][0]),cm2emu(line_ed[i+1][1]))
            
    def save(self):
        self.prs.save(self.file_name)
        
    def draw_shape(self,shape_id,x,y,dx,dy,adjs=[],rotation=0,fill_color = None,line_color=RGBColor(0, 0, 0),line_width=1):
        #1, Make Shape
        shape = self.slide.shapes.add_shape(shape_id,x,y,dx,dy)
        #2, Set Fill
        if type(fill_color) == type(None):
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        #3, Set Line
        if type(line_color) == type(None):
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(line_width)
        #Set Adjustments
        for (i,v) in adjs:
            shape.adjustments[i] = v
        shape.rotation = rotation
        shape.shadow.inherit = False

    def draw_textbox(self,x,y,dx,dy,margin,text,color,align="center"):
        tf = self.slide.shapes.add_textbox(x,y,dx,dy).text_frame
        tf.margin_bottom = cm2emu(0.0)
        tf.margin_top = cm2emu(0.0)
        tf.margin_left = cm2emu(margin)
        tf.margin_right = cm2emu(margin)
        para = tf.paragraphs[0]
        para.font.size = Pt(self.char_size)
        para.font.bold = True
        para.font.name = "Courier New"
        if align == "center":
            para.alignment = PP_ALIGN.CENTER
        else:
            para.alignment = PP_ALIGN.LEFT
        for i in range(len(text)):
            para.add_run()
            para.runs[i].text = text[i]
            if text[i] in self.paras["res_color"].keys():
                para.runs[i].font.color.rgb = RGBColor(self.paras["res_color"][text[i]][0],
                                                       self.paras["res_color"][text[i]][1],
                                                       self.paras["res_color"][text[i]][2])
            
            else:
                para.runs[i].font.color.rgb = RGBColor(0,0,0)

    
    def draw_line(self,stx,sty,edx,edy,line_color=RGBColor(0, 0, 0)):
        line = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,stx,sty,edx,edy)
        line.line.fill.solid()
        line.line.fill.fore_color.rgb = line_color
        line.line.width = 12700
        line.shadow.inherit = False
        
    def draw_resids(self,sh,stx,sty):
        if sh.st_id:
            self.draw_textbox(cm2emu(stx),cm2emu(sty-self.shape_height),
                              cm2emu(sh.length)/len(sh.resnames)*len(sh.resids[0]),
                              cm2emu(self.shape_height),sh.margin,
                                     sh.resids[0],"B",align="left")
        if sh.ed_id and len(sh.resnames) > len(sh.resids[0]) + 1:
            self.draw_textbox(cm2emu(stx+(len(sh.resids)-1)*self.char_length),cm2emu(sty-self.shape_height),
                              cm2emu(sh.length)/len(sh.resnames)*len(sh.resids[-1]),
                              cm2emu(self.shape_height),sh.margin,
                                     sh.resids[-1],"B",align="left")
        
        
    def run(self):
        self.make_PPT()
        self.make_rows()
        self.draw()
        try:
            self.save()
        except:
            self.error = "PPTX file could not be saved!\nCheck whether the file with the same name is open,\nor you have no permission to the output directory."

class Rows(list):
    def __init__(self,ss):
        self.update(ss)
        
    #Add a new row starting with a shape "ss"
    def update(self,ss):
        self.append(Row(ss))
        self.current_row = self[-1]


class Row(list):
    def __init__(self,ss):
        self.update(ss)
        
    #Add a new shape ss
    def update(self,ss):
        self.append(Shape(ss))
        self.current_shape = self[-1]
        
    def get_length(self,char_length):
        length = 0
        for s in self:
            if s.length > 0:
                length += s.length
        if len(self) > 0:
            length += (len(self) - 1) * char_length
        return length  
    
    
    
class Shape():
    def __init__(self,ss,st_id=True,ed_id=True):
        self.ss = ss
        self.resnames = ""
        self.colors = []
        self.resids = []
        self.length = 0
        self.margin = 0
        self.st_id = st_id
        self.ed_id = ed_id
        
        
    def update(self,resname,resid,paras,color=(0,0,0),char_length=0.13):
        self.resnames += resname
        self.resids.append(resid)
        self.colors.append(RGBColor(color[0],color[1],color[2]))
        self.update_length(paras=paras,char_length=char_length)
        
    def update_length(self,paras,char_length=0.13):
        text_length = len(self.resnames) * char_length
        if self.ss == "":
            shape_length = paras["kl_oval"] * text_length
            margin = 0.5 * (shape_length - text_length)
        elif self.ss == "S":
            margin = paras["margin_pentagon"]
            shape_length = text_length + 2 * margin
        elif self.ss == "H":
            margin = paras["margin_can"]
            shape_length = text_length + 2 * margin
        self.length = shape_length
        self.margin = margin
        



